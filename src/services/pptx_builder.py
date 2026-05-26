"""
Fixed-template PPTX generation for IC Memo.
7 slides with consistent layout — no AI-generated design.
"""
import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.models.agent_output import Vote
from src.models.memo import ICMemo

logger = logging.getLogger(__name__)

# ── Layout constants ──────────────────────────────────────────────────────────

_W = Inches(10)
_H = Inches(7.5)

_C = {
    "navy":    RGBColor(0x1F, 0x38, 0x64),
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "light":   RGBColor(0xF2, 0xF2, 0xF2),
    "dark":    RGBColor(0x33, 0x33, 0x33),
    "gray":    RGBColor(0x76, 0x76, 0x76),
    "green":   RGBColor(0x1E, 0x85, 0x44),
    "red":     RGBColor(0xC0, 0x39, 0x2B),
    "orange":  RGBColor(0xE6, 0x7E, 0x22),
    "gold":    RGBColor(0xC9, 0xA8, 0x4C),
}

_VOTE_COLOR = {Vote.APPROVE: _C["green"], Vote.REJECT: _C["red"], Vote.CONDITIONAL: _C["orange"]}
_VOTE_TEXT  = {Vote.APPROVE: "APPROVE",   Vote.REJECT: "REJECT",  Vote.CONDITIONAL: "CONDITIONAL"}
_SEVERITY_COLOR = {"High": _C["red"], "Medium": _C["orange"], "Low": _C["green"]}


# ── Public entry point ────────────────────────────────────────────────────────

def export_pptx(memo: ICMemo, output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    blank = prs.slide_layouts[6]  # blank layout

    _slide_cover(prs.slides.add_slide(blank), memo)
    _slide_overview(prs.slides.add_slide(blank), memo)
    _slide_financials(prs.slides.add_slide(blank), memo)
    _slide_thesis(prs.slides.add_slide(blank), memo)
    _slide_risks(prs.slides.add_slide(blank), memo)
    _slide_shock(prs.slides.add_slide(blank), memo)
    _slide_recommendation(prs.slides.add_slide(blank), memo)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("PPTX exported: %s (%d slides)", output_path, len(prs.slides))
    return output_path


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_cover(slide, memo: ICMemo) -> None:
    h = memo.header
    _fill_bg(slide, _C["navy"])
    _text_box(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.2),
              h.company_name, 40, bold=True, color=_C["white"], align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(1), Inches(3.4), Inches(8), Inches(0.6),
              "Investment Committee Memo", 18, color=_C["gold"], align=PP_ALIGN.CENTER)
    detail = (
        f"{h.industry}  ·  {h.deal_stage or 'N/A'}  ·  "
        f"${h.investment_amount_usd_m:.1f}M  ·  {h.prepared_date}"
    )
    _text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
              detail, 13, color=_C["light"], align=PP_ALIGN.CENTER)
    decision = memo.recommendation.decision
    _text_box(slide, Inches(3.5), Inches(5.3), Inches(3), Inches(0.7),
              _VOTE_TEXT[decision], 22, bold=True,
              color=_VOTE_COLOR[decision], align=PP_ALIGN.CENTER)


def _slide_overview(slide, memo: ICMemo) -> None:
    _add_title_bar(slide, "1. Investment Overview")
    items = memo.overview.key_facts or ["No data collected yet."]
    _bullet_list(slide, Inches(0.4), Inches(1.4), Inches(9.2), Inches(5.4), items, size=14)
    _footer(slide, memo.header.company_name, "1")


def _slide_financials(slide, memo: ICMemo) -> None:
    _add_title_bar(slide, "2. Financial Analysis")
    f = memo.financials
    irr_str = f"{f.base_irr_pct:.1f}%" if f.base_irr_pct is not None else "—"

    # Metric cards (top row)
    metrics = [
        ("Base IRR",      irr_str),
        ("Base MOIC",     f"{f.base_moic:.2f}x"),
        ("Exit Value",    f"${f.exit_value_usd_m:.1f}M"),
        ("Exit Multiple", f"{f.exit_multiple:.1f}x Rev"),
    ]
    for i, (label, val) in enumerate(metrics):
        x = Inches(0.3 + i * 2.45)
        _metric_card(slide, x, Inches(1.35), Inches(2.2), label, val)

    # Findings + concerns
    col_items = [f"▸ {item}" for item in f.findings] + \
                [f"⚠ {item}" for item in f.concerns]
    _bullet_list(slide, Inches(0.4), Inches(2.9), Inches(9.2), Inches(3.8),
                 col_items or ["Pending financial analysis."], size=12)
    _footer(slide, memo.header.company_name, "2")


def _slide_thesis(slide, memo: ICMemo) -> None:
    _add_title_bar(slide, "3. Investment Thesis")
    t = memo.thesis
    # Two-column: bull | bear
    _text_box(slide, Inches(0.3), Inches(1.35), Inches(4.6), Inches(0.4),
              "Bull Case", 13, bold=True, color=_C["green"])
    _bullet_list(slide, Inches(0.3), Inches(1.8), Inches(4.6), Inches(4.8),
                 t.bull_points or ["—"], size=12)
    _text_box(slide, Inches(5.1), Inches(1.35), Inches(4.6), Inches(0.4),
              "Bear Case", 13, bold=True, color=_C["red"])
    _bullet_list(slide, Inches(5.1), Inches(1.8), Inches(4.6), Inches(4.8),
                 t.bear_points or ["—"], size=12)
    _footer(slide, memo.header.company_name, "3")


def _slide_risks(slide, memo: ICMemo) -> None:
    _add_title_bar(slide, "4. Key Risks")
    risks = memo.risks.risks
    y = Inches(1.35)
    row_h = Inches(0.55)
    # Header row
    _text_box(slide, Inches(0.3), y, Inches(0.5), row_h, "#", 11, bold=True, color=_C["white"], bg=_C["navy"])
    _text_box(slide, Inches(0.85), y, Inches(6.8), row_h, "Risk Description", 11, bold=True, color=_C["white"], bg=_C["navy"])
    _text_box(slide, Inches(7.7), y, Inches(2.0), row_h, "Severity", 11, bold=True, color=_C["white"], bg=_C["navy"])
    for i, risk in enumerate(risks[:8]):
        y += row_h
        bg = _C["light"] if i % 2 == 0 else _C["white"]
        sev_color = _SEVERITY_COLOR.get(risk.severity, _C["gray"])
        _text_box(slide, Inches(0.3),  y, Inches(0.5),  row_h, str(i + 1),      11, bg=bg)
        _text_box(slide, Inches(0.85), y, Inches(6.8),  row_h, risk.description, 11, bg=bg)
        _text_box(slide, Inches(7.7),  y, Inches(2.0),  row_h, risk.severity,    11, bold=True, color=sev_color, bg=bg)
    _footer(slide, memo.header.company_name, "4")


def _slide_shock(slide, memo: ICMemo) -> None:
    _add_title_bar(slide, "5. Shock Simulation Summary")
    s = memo.shock_summary
    base_str = (
        f"Base Case: IRR {s.base_irr_pct:.1f}%  |  MOIC {s.base_moic:.2f}x"
        if s.base_irr_pct is not None else "Base Case: N/A"
    )
    _text_box(slide, Inches(0.4), Inches(1.35), Inches(9.2), Inches(0.4),
              base_str, 12, bold=True, color=_C["navy"])
    if s.top_driver:
        _text_box(slide, Inches(0.4), Inches(1.8), Inches(9.2), Inches(0.35),
                  f"Top Sensitivity Driver: {s.top_driver}  (±{s.top_driver_swing_pp:.1f}pp IRR swing)",
                  11, color=_C["gray"])
    # Table
    y = Inches(2.25)
    row_h = Inches(0.48)
    _text_box(slide, Inches(0.3), y, Inches(3.8), row_h, "Scenario",    11, bold=True, color=_C["white"], bg=_C["navy"])
    _text_box(slide, Inches(4.2), y, Inches(1.8), row_h, "IRR",         11, bold=True, color=_C["white"], bg=_C["navy"])
    _text_box(slide, Inches(6.1), y, Inches(1.5), row_h, "MOIC",        11, bold=True, color=_C["white"], bg=_C["navy"])
    _text_box(slide, Inches(7.7), y, Inches(2.0), row_h, "ΔIRR vs Base",11, bold=True, color=_C["white"], bg=_C["navy"])
    for i, row in enumerate(s.scenarios[:7]):
        y += row_h
        bg = _C["light"] if i % 2 == 0 else _C["white"]
        irr_str = f"{row.irr_pct:.1f}%" if row.irr_pct is not None else "—"
        delta = f"{row.delta_irr_pp:+.1f}pp" if row.delta_irr_pp is not None else "—"
        d_color = _C["red"] if (row.delta_irr_pp or 0) < 0 else _C["green"]
        _text_box(slide, Inches(0.3), y, Inches(3.8), row_h, row.label,     11, bg=bg)
        _text_box(slide, Inches(4.2), y, Inches(1.8), row_h, irr_str,       11, bg=bg)
        _text_box(slide, Inches(6.1), y, Inches(1.5), row_h, f"{row.moic:.2f}x", 11, bg=bg)
        _text_box(slide, Inches(7.7), y, Inches(2.0), row_h, delta,         11, bold=True, color=d_color, bg=bg)
    _footer(slide, memo.header.company_name, "5")


def _slide_recommendation(slide, memo: ICMemo) -> None:
    _add_title_bar(slide, "6. Final Recommendation")
    r = memo.recommendation
    decision_color = _VOTE_COLOR[r.decision]

    # Decision badge
    _text_box(slide, Inches(0.4), Inches(1.35), Inches(3.0), Inches(0.9),
              _VOTE_TEXT[r.decision], 28, bold=True, color=decision_color, align=PP_ALIGN.CENTER)
    quorum = "Quorum: Met" if r.quorum_met else "Quorum: Not Met"
    _text_box(slide, Inches(0.4), Inches(2.35), Inches(3.0), Inches(0.4),
              quorum, 11, color=_C["gray"])

    # Vote tally
    _text_box(slide, Inches(3.8), Inches(1.35), Inches(5.8), Inches(0.35),
              "Vote Tally", 12, bold=True, color=_C["navy"])
    tally_lines = [f"{agent}: {vote}" for agent, vote in r.vote_tally.items()] or ["—"]
    _bullet_list(slide, Inches(3.8), Inches(1.75), Inches(5.8), Inches(1.5), tally_lines, size=11)

    # Conditions
    if r.conditions:
        _text_box(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.35),
                  "Conditions", 12, bold=True, color=_C["orange"])
        _bullet_list(slide, Inches(0.4), Inches(3.4), Inches(9.2), Inches(1.2), r.conditions, size=11)

    # Rationale
    _text_box(slide, Inches(0.4), Inches(4.8), Inches(9.2), Inches(0.35),
              "Rationale", 12, bold=True, color=_C["navy"])
    _text_box(slide, Inches(0.4), Inches(5.2), Inches(9.2), Inches(1.6),
              r.rationale or "Pending committee rationale.", 11, color=_C["dark"], wrap=True)
    _footer(slide, memo.header.company_name, "6")


# ── Low-level drawing helpers ─────────────────────────────────────────────────

def _fill_bg(slide, color: RGBColor) -> None:
    from pptx.util import Emu
    sp = slide.shapes.add_shape(1, 0, 0, _W, _H)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()


def _add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(1, 0, 0, _W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C["navy"]
    bar.line.fill.background()
    _text_box(slide, Inches(0.3), Inches(0.18), Inches(9.4), Inches(0.75),
              title, 20, bold=True, color=_C["white"])


def _metric_card(slide, x, y, w, label: str, value: str) -> None:
    card = slide.shapes.add_shape(1, x, y, w, Inches(1.35))
    card.fill.solid()
    card.fill.fore_color.rgb = _C["light"]
    card.line.color.rgb = _C["navy"]
    _text_box(slide, x, y + Inches(0.12), w, Inches(0.45), label, 10, color=_C["gray"], align=PP_ALIGN.CENTER)
    _text_box(slide, x, y + Inches(0.55), w, Inches(0.65), value, 20, bold=True, color=_C["navy"], align=PP_ALIGN.CENTER)


def _text_box(
    slide, left, top, width, height, text: str, size: int,
    bold: bool = False, color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    bg: RGBColor | None = None,
    wrap: bool = False,
) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
        txb.line.fill.background()
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or _C["dark"]


def _bullet_list(slide, left, top, width, height, items: list[str], size: int = 12) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = _C["dark"]
        p.space_after = Pt(3)


def _footer(slide, company: str, page: str) -> None:
    _text_box(slide, Inches(0.3), Inches(7.1), Inches(6), Inches(0.3),
              f"Confidential — {company}", 9, color=_C["gray"])
    _text_box(slide, Inches(9.0), Inches(7.1), Inches(0.7), Inches(0.3),
              page, 9, color=_C["gray"], align=PP_ALIGN.RIGHT)
